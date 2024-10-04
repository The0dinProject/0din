import os
import sys
import subprocess
import zipfile
import tarfile
from PIL import Image, ImageDraw, ImageFont
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import ebooklib
from ebooklib import epub
import matplotlib.pyplot as plt
from pydub import AudioSegment
import tempfile

def generate_image_preview(input_file, output_file):
    try:
        # Identify the file type based on extension
        file_ext = os.path.splitext(input_file)[1].lower()

        if file_ext in ['.jpg', '.jpeg', '.png', '.bmp', '.gif', '.tiff']:
            # Process image files
            process_image(input_file, output_file)
        elif file_ext in ['.mp4', '.mkv', '.avi', '.mov', '.webm']:
            # Process video files
            process_video(input_file, output_file)
        elif file_ext in ['.mp3', '.wav', '.ogg', '.flac']:
            # Process audio files
            process_audio(input_file, output_file)
        elif file_ext == '.pdf':
            # Process PDF files
            process_pdf(input_file, output_file)
        elif file_ext == '.docx':
            # Process DOCX files
            process_docx(input_file, output_file)
        elif file_ext == '.pptx':
            # Process PPTX files
            process_pptx(input_file, output_file)
        elif file_ext == '.epub':
            # Process EPUB files
            process_epub(input_file, output_file)
        elif file_ext in ['.txt', '.md', '.py', '.html', '.css', '.js']:
            # Process text and code files
            process_text(input_file, output_file)
        elif file_ext in ['.zip', '.tar', '.gz']:
            # Process archive files
            process_archive(input_file, output_file)
        else:
            # Fallback for unsupported file types
            process_generic_placeholder(output_file)
    except Exception as e:
        print(f"Error processing file: {e}")

# ------------------- Helper functions for each format ------------------- #

def process_image(input_file, output_file):
    """Downscale image and save as webp"""
    try:
        with Image.open(input_file) as img:
            img.thumbnail((512, 512))  # Resize to 512x512 max
            img.save(output_file, "WEBP")
        print(f"Image preview saved at {output_file}")
    except Exception as e:
        print(f"Failed to process image: {e}")

def process_video(input_file, output_file):
    """Generate a video thumbnail using ffmpegthumbnailer"""
    try:
        command = ['ffmpegthumbnailer', '-i', input_file, '-o', output_file, '-s', '512', '-f']
        subprocess.run(command, check=True)
        print(f"Video preview saved at {output_file}")
    except Exception as e:
        print(f"Failed to process video: {e}")

def process_audio(input_file, output_file):
    """Generate a waveform image from an audio file"""
    try:
        # Load audio file
        audio = AudioSegment.from_file(input_file)
        data = audio.get_array_of_samples()

        # Plot waveform
        plt.figure(figsize=(8, 4))
        plt.plot(data[:10000])  # Only plot first 10k samples for preview
        plt.axis('off')
        
        # Save the waveform as an image
        plt.savefig(output_file, format="webp", bbox_inches='tight', pad_inches=0)
        plt.close()
        print(f"Audio preview saved at {output_file}")
    except Exception as e:
        print(f"Failed to process audio: {e}")

def process_pdf(input_file, output_file):
    """Generate a thumbnail from the first page of a PDF"""
    try:
        pdf_document = fitz.open(input_file)
        page = pdf_document.load_page(0)  # Get the first page
        pix = page.get_pixmap()
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        img.thumbnail((512, 512))
        img.save(output_file, "WEBP")
        print(f"PDF preview saved at {output_file}")
    except Exception as e:
        print(f"Failed to process PDF: {e}")

def process_docx(input_file, output_file):
    """Generate a thumbnail from the first page of a DOCX document"""
    try:
        doc = Document(input_file)
        if doc.paragraphs:
            text = doc.paragraphs[0].text
        else:
            text = "No content"
        
        # Create a blank image and draw text
        img = Image.new("RGB", (512, 512), (255, 255, 255))
        d = ImageDraw.Draw(img)
        d.text((10, 10), text[:200], fill=(0, 0, 0))  # Show the first 200 characters
        img.save(output_file, "WEBP")
        print(f"DOCX preview saved at {output_file}")
    except Exception as e:
        print(f"Failed to process DOCX: {e}")

def process_pptx(input_file, output_file):
    """Generate a thumbnail from the first slide of a PPTX presentation"""
    try:
        prs = Presentation(input_file)
        first_slide = prs.slides[0]
        
        # Create an image with the title text
        if first_slide.shapes.title:
            title = first_slide.shapes.title.text
        else:
            title = "No Title"
        
        img = Image.new("RGB", (512, 512), (255, 255, 255))
        d = ImageDraw.Draw(img)
        d.text((10, 10), title[:200], fill=(0, 0, 0))  # Show the first 200 characters
        img.save(output_file, "WEBP")
        print(f"PPTX preview saved at {output_file}")
    except Exception as e:
        print(f"Failed to process PPTX: {e}")

def process_epub(input_file, output_file):
    """Generate a thumbnail from an EPUB ebook"""
    try:
        book = epub.read_epub(input_file)
        cover = None

        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_COVER:
                cover = item.content
                break

        if cover:
            with open(tempfile.mktemp(suffix=".jpg"), 'wb') as f:
                f.write(cover)
                img = Image.open(f.name)
                img.thumbnail((512, 512))
                img.save(output_file, "WEBP")
        else:
            process_generic_placeholder(output_file)
        print(f"EPUB preview saved at {output_file}")
    except Exception as e:
        print(f"Failed to process EPUB: {e}")

def process_text(input_file, output_file):
    """Generate a preview from text or code file"""
    try:
        with open(input_file, 'r') as f:
            text = f.read(200)  # Read the first 200 characters
        
        img = Image.new("RGB", (512, 512), (255, 255, 255))
        d = ImageDraw.Draw(img)
        d.text((10, 10), text, fill=(0, 0, 0))
        img.save(output_file, "WEBP")
        print(f"Text preview saved at {output_file}")
    except Exception as e:
        print(f"Failed to process text file: {e}")

def process_archive(input_file, output_file):
    """Generate a preview of archive contents"""
    try:
        if zipfile.is_zipfile(input_file):
            with zipfile.ZipFile(input_file, 'r') as archive:
                file_list = archive.namelist()[:10]  # Show first 10 files
        elif tarfile.is_tarfile(input_file):
            with tarfile.open(input_file, 'r') as archive:
                file_list = [tarinfo.name for tarinfo in archive.getmembers()][:10]
        else:
            file_list = []

        img = Image.new("RGB", (512, 512), (255, 255, 255))
        d = ImageDraw.Draw(img)
        d.text((10, 10), "\n".join(file_list), fill=(0, 0, 0))  # Display file names
        img.save(output_file, "WEBP")
        print(f"Archive preview saved at {output_file}")
    except Exception as e:
        print(f"Failed to process archive: {e}")

def process_generic_placeholder(output_file):
    """Generate a placeholder preview for unsupported file types"""
    try:
        img = Image.new("RGB", (512, 512), (200, 200, 200))
        d = ImageDraw.Draw(img)
        d.text((100, 250), "Preview Not Available", fill=(0, 0, 0))
        img.save(output_file, "WEBP")
        print(f"Generic placeholder saved at {output_file}")
    except Exception as e:
        print(f"Failed to create placeholder: {e}")

# ---------------------------- Main Execution ---------------------------- #

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python generate_preview.py <input_file> <output_file>")
        sys.exit(1)

    input_file = sys.argv[1]
    output_file = sys.argv[2]

    generate_image_preview(input_file, output_file)

